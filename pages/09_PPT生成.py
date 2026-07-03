import streamlit as st
import json
import os
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from data_tools import add_resource, RESOURCE_FILE_DIR, get_all_class
import requests

st.set_page_config(page_title="PPT生成", layout="centered", initial_sidebar_state="collapsed")

# 教师权限校验
if "teacher_verified" not in st.session_state or not st.session_state["teacher_verified"]:
    st.warning("请先通过教师入口验证")
    st.page_link("pages/00_教师入口首页.py", label="前往验证", use_container_width=True)
    st.stop()

st.markdown("<h1 style='text-align:center;'>教学PPT一键生成</h1>", unsafe_allow_html=True)
st.divider()

# ========== 安全读取密钥 ==========
API_KEY = ""
API_BASE = ""
MODEL_NAME = ""

try:
    if "DEEPSEEK_API_KEY" in st.secrets:
        API_KEY = st.secrets["DEEPSEEK_API_KEY"]
        API_BASE = "https://api.deepseek.com/v1"
        MODEL_NAME = "deepseek-chat"
    elif "DASHSCOPE_API_KEY" in st.secrets:
        API_KEY = st.secrets["DASHSCOPE_API_KEY"]
        API_BASE = "https://dashscope.aliyuncs.com/compatible-mode/v1"
        MODEL_NAME = "qwen-plus"
except Exception:
    pass

# ========== 两套基础教学模板 ==========
THEMES = {
    "简约商务蓝（通用）": {
        "bg_start": RGBColor(240, 246, 255),
        "bg_end": RGBColor(255, 255, 255),
        "main_color": RGBColor(22, 93, 255),
        "sub_color": RGBColor(100, 149, 237),
        "text_title": RGBColor(25, 50, 100),
        "text_body": RGBColor(60, 60, 60)
    },
    "国风靛蓝朱砂（民族文化）": {
        "bg_start": RGBColor(245, 240, 235),
        "bg_end": RGBColor(255, 253, 250),
        "main_color": RGBColor(25, 55, 110),
        "sub_color": RGBColor(180, 30, 30),
        "text_title": RGBColor(20, 40, 80),
        "text_body": RGBColor(50, 50, 50)
    }
}


# ========== 1. 快速生成：按主题生成大纲 ==========
def generate_outline_by_topic(topic, subject, page_count):
    prompt = f"""
    请为中职{subject}学科生成一份主题为《{topic}》的教学PPT大纲，正文共{page_count}页。
    要求：
    1. 内容符合中职教学难度，逻辑清晰，适合课堂授课
    2. 每页包含标题和3-5个核心要点，要点简洁凝练
    3. 最后一页为课堂小结+作业布置
    4. 严格返回JSON格式，不要任何额外文字，格式如下：
    {{
        "main_title": "完整课件标题",
        "slides": [
            {{"title": "页面标题", "points": ["要点1", "要点2", "要点3"]}}
        ]
    }}
    """
    return call_ai_api(prompt)


# ========== 2. 教案联动：从教案文本生成大纲 ==========
def generate_outline_by_lesson_plan(lesson_plan, subject, page_count):
    prompt = f"""
    以下是一份中职{subject}学科的教案原文，请你基于教案内容，提炼生成一份{page_count}页的教学PPT大纲。
    要求：
    1. 严格基于教案内容提炼，不新增无关知识点，贴合授课流程
    2. 按「课程导入-知识讲解-案例演示-课堂小结-作业布置」的教学逻辑组织页面
    3. 每页包含标题和3-5个核心要点，要点简洁凝练，适合投影展示
    4. 严格返回JSON格式，不要任何额外文字，格式如下：
    {{
        "main_title": "完整课件标题",
        "slides": [
            {{"title": "页面标题", "points": ["要点1", "要点2", "要点3"]}}
        ]
    }}

    教案原文：
    {lesson_plan}
    """
    return call_ai_api(prompt)


# ========== 通用AI调用函数 ==========
def call_ai_api(prompt):
    try:
        headers = {
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json"
        }
        data = {
            "model": MODEL_NAME,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.7,
            "response_format": {"type": "json_object"}
        }
        res = requests.post(f"{API_BASE}/chat/completions", headers=headers, json=data, timeout=120)
        res.raise_for_status()
        content = res.json()["choices"][0]["message"]["content"]
        return json.loads(content)
    except Exception as e:
        st.error(f"AI生成失败：{str(e)}")
        return None


# ========== PPT渲染引擎（两套模式共用） ==========
def set_slide_bg(slide, theme):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops[0].color.rgb = theme["bg_start"]
    fill.gradient_stops[1].color.rgb = theme["bg_end"]


def add_decoration(slide, theme):
    left, top, width, height = 0, 0, Inches(13.33), Inches(0.08)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["main_color"]
    shape.line.fill.background()

    left, top, width, height = Inches(0.5), Inches(7.2), Inches(12.33), Emu(1)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = theme["sub_color"]
    line.line.fill.background()


def add_page_num(slide, num, total, theme):
    txBox = slide.shapes.add_textbox(Inches(12), Inches(7.3), Inches(1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = theme["text_body"]
    p.alignment = PP_ALIGN.RIGHT


def create_ppt(outline_data, subject, theme_name):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    theme = THEMES[theme_name]
    slides_content = outline_data["slides"]
    total_pages = len(slides_content) + 3

    # 封面页
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(cover, theme)
    shape = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3), Inches(13.33), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["main_color"]
    shape.line.fill.background()

    title_box = cover.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.33), Inches(1.2))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = outline_data["main_title"]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "微软雅黑"

    sub_box = cover.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(0.6))
    sub_tf = sub_box.text_frame
    sub_tf.word_wrap = True
    p = sub_tf.paragraphs[0]
    p.text = f"{subject} · 教学课件"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(240, 240, 240)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "微软雅黑"

    # 目录页
    catalog = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(catalog, theme)
    add_decoration(catalog, theme)
    add_page_num(catalog, 2, total_pages, theme)

    cat_title = catalog.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.33), Inches(0.8))
    cat_tf = cat_title.text_frame
    p = cat_tf.paragraphs[0]
    p.text = "目 录"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = theme["text_title"]
    p.font.name = "微软雅黑"

    line = catalog.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.6), Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = theme["sub_color"]
    line.line.fill.background()

    list_top = 2.2
    for idx, slide in enumerate(slides_content, 1):
        circle = catalog.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(list_top), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme["main_color"]
        circle.line.fill.background()
        circle_tf = circle.text_frame
        circle_tf.paragraphs[0].text = str(idx)
        circle_tf.paragraphs[0].font.size = Pt(18)
        circle_tf.paragraphs[0].font.bold = True
        circle_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        circle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        text_box = catalog.shapes.add_textbox(Inches(2), Inches(list_top + 0.05), Inches(10), Inches(0.5))
        text_tf = text_box.text_frame
        p = text_tf.paragraphs[0]
        p.text = slide["title"]
        p.font.size = Pt(22)
        p.font.color.rgb = theme["text_body"]
        p.font.name = "微软雅黑"
        list_top += 0.8

    # 正文页
    for idx, slide_data in enumerate(slides_content):
        page_num = idx + 3
        page = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(page, theme)
        add_decoration(page, theme)
        add_page_num(page, page_num, total_pages, theme)

        title_box = page.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.73), Inches(0.8))
        title_tf = title_box.text_frame
        p = title_tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = theme["text_title"]
        p.font.name = "微软雅黑"

        line = page.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(1), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = theme["sub_color"]
        line.line.fill.background()

        body_top = 1.9
        for point in slide_data["points"]:
            dot = page.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(body_top + 0.12), Inches(0.15),
                                        Inches(0.15))
            dot.fill.solid()
            dot.fill.fore_color.rgb = theme["main_color"]
            dot.line.fill.background()

            point_box = page.shapes.add_textbox(Inches(1.6), Inches(body_top), Inches(11), Inches(0.6))
            point_tf = point_box.text_frame
            point_tf.word_wrap = True
            p = point_tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = theme["text_body"]
            p.font.name = "微软雅黑"
            p.space_after = Pt(8)
            body_top += 0.7

    # 结尾页
    end = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(end, theme)
    add_decoration(end, theme)

    end_box = end.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(1.5))
    end_tf = end_box.text_frame
    p = end_tf.paragraphs[0]
    p.text = "感谢观看"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = theme["main_color"]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "微软雅黑"

    sub_box = end.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(0.6))
    sub_tf = sub_box.text_frame
    p = sub_tf.paragraphs[0]
    p.text = "欢迎批评指正"
    p.font.size = Pt(24)
    p.font.color.rgb = theme["text_body"]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "微软雅黑"

    # 保存文件
    safe_title = outline_data["main_title"].replace("/", "_").replace("\\", "_")
    file_name = f"PPT_{safe_title}.pptx"
    save_path = os.path.join(RESOURCE_FILE_DIR, file_name)
    prs.save(save_path)
    return save_path, file_name


# ========== 页面交互 ==========
if not API_KEY:
    st.warning("未检测到AI密钥，请先在项目根目录的 .streamlit/secrets.toml 中配置API密钥")
    st.page_link("pages/00_教师入口首页.py", label="返回教师首页", use_container_width=True)
    st.stop()

# 模式切换
gen_mode = st.radio("选择生成模式", ["快速生成模式", "教案联动模式"], horizontal=True)
st.divider()

# 公共参数：模板风格
style_select = st.selectbox("选择模板风格", list(THEMES.keys()))

# 模式1：快速生成
if gen_mode == "快速生成模式":
    col1, col2 = st.columns(2)
    with col1:
        ppt_topic = st.text_input("授课主题", value="轴对称图形")
        ppt_subject = st.text_input("学科", value="数学基础模块")
    with col2:
        page_count = st.number_input("正文页数", min_value=3, max_value=20, value=8, step=1)

# 模式2：教案联动
else:
    ppt_subject = st.text_input("学科", value="数学基础模块")
    lesson_plan_text = st.text_area(
        "粘贴教案内容（可直接从备课助手复制）",
        height=260,
        placeholder="在此粘贴完整教案文本，AI将自动提炼内容生成PPT..."
    )
    page_count = st.number_input("生成正文页数", min_value=3, max_value=20, value=10, step=1)

st.divider()

# 生成按钮
if st.button("生成PPT课件", type="primary", use_container_width=True):
    outline = None
    if gen_mode == "快速生成模式":
        if not ppt_topic.strip():
            st.warning("请填写授课主题")
        else:
            with st.spinner("AI正在生成课件大纲，请稍候..."):
                outline = generate_outline_by_topic(ppt_topic.strip(), ppt_subject.strip(), page_count)
    else:
        if not lesson_plan_text.strip():
            st.warning("请粘贴教案内容")
        else:
            with st.spinner("AI正在解析教案并生成大纲，请稍候..."):
                outline = generate_outline_by_lesson_plan(lesson_plan_text.strip(), ppt_subject.strip(), page_count)

    if outline:
        with st.spinner("正在生成PPT文件..."):
            try:
                path, name = create_ppt(outline, ppt_subject.strip(), style_select)
                st.session_state["last_ppt"] = {"path": path, "name": name, "topic": outline["main_title"]}
                st.rerun()
            except Exception as e:
                st.error(f"PPT生成失败：{str(e)}")

# 生成结果区
if "last_ppt" in st.session_state and os.path.exists(st.session_state["last_ppt"]["path"]):
    st.divider()
    st.subheader("生成结果")
    st.success(f"课件已生成：{st.session_state['last_ppt']['name']}")

    col_dl, col_save = st.columns(2)
    with col_dl:
        with open(st.session_state["last_ppt"]["path"], "rb") as f:
            st.download_button(
                label="📥 下载PPT文件",
                data=f.read(),
                file_name=st.session_state["last_ppt"]["name"],
                use_container_width=True
            )
    with col_save:
        class_list = get_all_class()
        assign_cls = st.multiselect("分配班级", class_list, default=class_list, key="ppt_assign")
        if st.button("保存到资源管理", use_container_width=True):
            if len(assign_cls) == 0:
                st.error("请至少选择一个班级")
            else:
                res_id = add_resource(
                    f"{st.session_state['last_ppt']['topic']}课件",
                    "课件文档",
                    st.session_state["last_ppt"]["name"],
                    assign_cls
                )
                st.success(f"已保存到资源库，资源ID：{res_id}")

st.divider()
st.page_link("pages/00_教师入口首页.py", label="← 返回教师首页", use_container_width=True)